from pathlib import Path
from typing import Annotated
import os
import tempfile
import zipfile
import subprocess
import shutil

import fitz  # PyMuPDF
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, PlainTextResponse

from pdf2docx import Converter
import pdfplumber
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


app = FastAPI(title="FileNova Conversion API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
    "http://localhost:5173",
    "http://127.0.0.1:5173",
    "http://localhost:5174",
    "http://127.0.0.1:5174",
    "https://filenova-web.onrender.com",
],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["X-Output-Filename"],
)


@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "FileNova Conversion API",
    }


# ============================================================
# PDF → WORD
# DO NOT CHANGE THIS WORKING CONVERSION
# ============================================================

def convert_pdf_to_docx(pdf_bytes: bytes):
    """
    Convert PDF to a real editable DOCX using pdf2docx.
    """
    pdf_path = None
    docx_path = None
    converter = None

    try:
        with tempfile.NamedTemporaryFile(
            suffix=".pdf",
            delete=False,
        ) as pdf_file:
            pdf_file.write(pdf_bytes)
            pdf_path = pdf_file.name

        with tempfile.NamedTemporaryFile(
            suffix=".docx",
            delete=False,
        ) as docx_file:
            docx_path = docx_file.name

        converter = Converter(pdf_path)

        converter.convert(
            docx_path,
            start=0,
            end=None,
        )

        converter.close()
        converter = None

        with open(docx_path, "rb") as generated:
            output = generated.read()

        if not output:
            raise RuntimeError(
                "The conversion produced an empty DOCX file."
            )

        return output

    finally:
        if converter is not None:
            try:
                converter.close()
            except Exception:
                pass

        if pdf_path:
            try:
                os.remove(pdf_path)
            except OSError:
                pass

        if docx_path:
            try:
                os.remove(docx_path)
            except OSError:
                pass


# ============================================================
# PDF → JPG
# ============================================================

def convert_pdf_to_jpg(pdf_bytes: bytes, original_filename: str):
    """
    Convert every PDF page into a high-quality JPG.

    One page:
        returns one JPG

    Multiple pages:
        returns a ZIP containing all JPG pages
    """

    pdf = None
    temp_dir = None

    try:
        pdf = fitz.open(
            stream=pdf_bytes,
            filetype="pdf",
        )

        if pdf.page_count == 0:
            raise RuntimeError(
                "The PDF does not contain any pages."
            )

        temp_dir = tempfile.mkdtemp(
            prefix="filenova_pdf_jpg_"
        )

        base_name = Path(
            original_filename
        ).stem

        jpg_files = []

        # 2.0 gives good quality without making files unnecessarily huge.
        scale = 2.0

        for page_number in range(pdf.page_count):
            page = pdf.load_page(page_number)

            matrix = fitz.Matrix(
                scale,
                scale,
            )

            pixmap = page.get_pixmap(
                matrix=matrix,
                alpha=False,
            )

            jpg_path = os.path.join(
                temp_dir,
                f"{base_name}_page_{page_number + 1}.jpg",
            )

            pixmap.save(
                jpg_path,
                output="jpeg",
            )

            jpg_files.append(jpg_path)

        pdf.close()
        pdf = None

        # ----------------------------------------------------
        # One page = direct JPG download
        # ----------------------------------------------------

        if len(jpg_files) == 1:
            jpg_path = jpg_files[0]

            with open(jpg_path, "rb") as image_file:
                output = image_file.read()

            output_filename = (
                f"{base_name}_page_1.jpg"
            )

            return (
                output,
                "image/jpeg",
                output_filename,
            )

        # ----------------------------------------------------
        # Multiple pages = ZIP download
        # ----------------------------------------------------

        zip_path = os.path.join(
            temp_dir,
            f"{base_name}_jpg.zip",
        )

        with zipfile.ZipFile(
            zip_path,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as zip_file:

            for jpg_path in jpg_files:
                zip_file.write(
                    jpg_path,
                    arcname=Path(jpg_path).name,
                )

        with open(zip_path, "rb") as zip_file:
            output = zip_file.read()

        output_filename = (
            f"{base_name}_jpg.zip"
        )

        return (
            output,
            "application/zip",
            output_filename,
        )

    finally:
        if pdf is not None:
            try:
                pdf.close()
            except Exception:
                pass

        if temp_dir:
            try:
                for filename in os.listdir(temp_dir):
                    filepath = os.path.join(
                        temp_dir,
                        filename,
                    )

                    try:
                        os.remove(filepath)
                    except OSError:
                        pass

                os.rmdir(temp_dir)

            except OSError:
                pass


# ============================================================
# WORD → PDF
# ============================================================

def convert_word_to_pdf(word_bytes: bytes, original_filename: str):
    """
    Convert a DOCX/DOC file to a real PDF using LibreOffice.

    LibreOffice is used as the conversion engine so the resulting
    PDF preserves the Word document's layout, tables, images and
    formatting much better than rebuilding the document manually.
    """
    temp_dir = tempfile.mkdtemp(prefix="filenova_word_pdf_")

    try:
        extension = Path(original_filename).suffix.lower()

        if extension not in {".docx", ".doc"}:
            raise ValueError(
                "Only DOCX and DOC files are supported."
            )

        input_filename = Path(original_filename).name
        input_path = os.path.join(temp_dir, input_filename)

        with open(input_path, "wb") as word_file:
            word_file.write(word_bytes)

        # Standard Windows LibreOffice installation path.
        # The x86 path is included as a fallback.
        libreoffice_candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]

        libreoffice = next(
            (
                path
                for path in libreoffice_candidates
                if os.path.exists(path)
            ),
            None,
        )

        # If LibreOffice is available through PATH, use that too.
        if libreoffice is None:
            libreoffice = shutil.which("soffice")

        if not libreoffice:
            raise RuntimeError(
                "LibreOffice was not found. "
                "Please install LibreOffice to use Word to PDF."
            )

        result = subprocess.run(
            [
                libreoffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                temp_dir,
                input_path,
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=120,
            creationflags=(
                subprocess.CREATE_NO_WINDOW
                if os.name == "nt"
                else 0
            ),
        )

        output_filename = (
            f"{Path(original_filename).stem}.pdf"
        )
        output_path = os.path.join(
            temp_dir,
            output_filename,
        )

        if result.returncode != 0:
            error_message = (
                result.stderr.strip()
                or result.stdout.strip()
                or "LibreOffice conversion failed."
            )
            raise RuntimeError(error_message)

        if not os.path.exists(output_path):
            raise RuntimeError(
                "LibreOffice completed but did not create the PDF."
            )

        with open(output_path, "rb") as pdf_file:
            output = pdf_file.read()

        if not output:
            raise RuntimeError(
                "The generated PDF is empty."
            )

        return output, output_filename

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True,
        )


# ============================================================
# PDF → EXCEL
# ============================================================

def convert_pdf_to_excel(pdf_bytes: bytes, original_filename: str):
    """
    Convert PDF tables into a real editable XLSX workbook.

    Each PDF page gets its own worksheet. Detected tables are written
    into the worksheet as editable Excel cells. If a page has no
    detectable table, its extracted text is written line-by-line so
    the user still receives useful editable content.
    """
    temp_dir = tempfile.mkdtemp(prefix="filenova_pdf_excel_")

    try:
        pdf_path = os.path.join(temp_dir, Path(original_filename).name)
        xlsx_filename = f"{Path(original_filename).stem}.xlsx"
        xlsx_path = os.path.join(temp_dir, xlsx_filename)

        with open(pdf_path, "wb") as pdf_file:
            pdf_file.write(pdf_bytes)

        workbook = Workbook()
        # Remove the default sheet; we create one sheet per PDF page.
        default_sheet = workbook.active
        workbook.remove(default_sheet)

        table_count = 0

        with pdfplumber.open(pdf_path) as pdf:
            if not pdf.pages:
                raise RuntimeError("The PDF does not contain any pages.")

            for page_number, page in enumerate(pdf.pages, start=1):
                # Excel sheet names have a 31-character limit.
                sheet = workbook.create_sheet(
                    title=f"Page {page_number}"[:31]
                )

                tables = page.extract_tables()

                if tables:
                    current_row = 1

                    for table in tables:
                        table_count += 1

                        for row in table:
                            if row is None:
                                continue

                            for column_number, value in enumerate(
                                row,
                                start=1,
                            ):
                                cell_value = (
                                    "" if value is None else str(value).strip()
                                )
                                sheet.cell(
                                    row=current_row,
                                    column=column_number,
                                    value=cell_value,
                                )

                            current_row += 1

                        # Leave one blank row between separate tables.
                        current_row += 1

                else:
                    # No table detected: preserve page text as editable
                    # rows instead of returning an empty workbook.
                    text_content = page.extract_text() or ""

                    if text_content.strip():
                        for row_number, line in enumerate(
                            text_content.splitlines(),
                            start=1,
                        ):
                            sheet.cell(
                                row=row_number,
                                column=1,
                                value=line,
                            )
                    else:
                        sheet.cell(
                            row=1,
                            column=1,
                            value="No extractable text or table found on this page.",
                        )

                # Make columns readable without changing the extracted data.
                for column_cells in sheet.columns:
                    max_length = 0
                    column_letter = column_cells[0].column_letter

                    for cell in column_cells:
                        if cell.value is not None:
                            max_length = max(
                                max_length,
                                len(str(cell.value)),
                            )

                    sheet.column_dimensions[column_letter].width = min(
                        max(max_length + 2, 10),
                        60,
                    )

        if table_count == 0:
            # The workbook still contains page text, but tell the user
            # clearly what happened in the server log.
            print(
                "PDF to Excel: no tables detected; "
                "page text was exported instead."
            )

        workbook.save(xlsx_path)

        with open(xlsx_path, "rb") as excel_file:
            output = excel_file.read()

        if not output:
            raise RuntimeError("The generated Excel file is empty.")

        return output, xlsx_filename

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True,
        )


# ============================================================
# EXCEL → PDF
# ============================================================

def convert_excel_to_pdf(excel_bytes: bytes, original_filename: str):
    """
    Convert an XLSX/XLS file to a real PDF using LibreOffice.
    LibreOffice handles workbook layout, sheets, tables and formatting.
    """
    temp_dir = tempfile.mkdtemp(prefix="filenova_excel_pdf_")

    try:
        extension = Path(original_filename).suffix.lower()

        if extension not in {".xlsx", ".xls"}:
            raise ValueError(
                "Only XLSX and XLS files are supported."
            )

        input_filename = Path(original_filename).name
        input_path = os.path.join(temp_dir, input_filename)

        with open(input_path, "wb") as excel_file:
            excel_file.write(excel_bytes)

        libreoffice_candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]

        libreoffice = next(
            (
                path
                for path in libreoffice_candidates
                if os.path.exists(path)
            ),
            None,
        )

        if libreoffice is None:
            libreoffice = shutil.which("soffice")

        if not libreoffice:
            raise RuntimeError(
                "LibreOffice was not found. "
                "Please install LibreOffice to use Excel to PDF."
            )

        result = subprocess.run(
            [
                libreoffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                temp_dir,
                input_path,
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=120,
            creationflags=(
                subprocess.CREATE_NO_WINDOW
                if os.name == "nt"
                else 0
            ),
        )

        output_filename = (
            f"{Path(original_filename).stem}.pdf"
        )
        output_path = os.path.join(
            temp_dir,
            output_filename,
        )

        if result.returncode != 0:
            error_message = (
                result.stderr.strip()
                or result.stdout.strip()
                or "LibreOffice conversion failed."
            )
            raise RuntimeError(error_message)

        if not os.path.exists(output_path):
            raise RuntimeError(
                "LibreOffice completed but did not create the PDF."
            )

        with open(output_path, "rb") as pdf_file:
            output = pdf_file.read()

        if not output:
            raise RuntimeError(
                "The generated PDF is empty."
            )

        return output, output_filename

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True,
        )


# ============================================================
# PDF → POWERPOINT
# ============================================================

def convert_pdf_to_powerpoint(pdf_bytes: bytes, original_filename: str):
    """
    Convert each PDF page into a PowerPoint slide.

    Each page is rendered at high resolution and placed as a
    full-slide image. This preserves the visual appearance of
    the original PDF reliably across pages.
    """
    temp_dir = tempfile.mkdtemp(prefix="filenova_pdf_pptx_")

    try:
        base_name = Path(original_filename).stem
        pptx_filename = f"{base_name}.pptx"
        pptx_path = os.path.join(temp_dir, pptx_filename)

        pdf = fitz.open(stream=pdf_bytes, filetype="pdf")

        if pdf.page_count == 0:
            pdf.close()
            raise RuntimeError("The PDF does not contain any pages.")

        # Use a standard 16:9 presentation size.
        prs = Presentation()
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)

        # Remove the default blank slide if one exists.
        if len(prs.slides) > 0:
            slide_ids = prs.slides._sldIdLst
            slide_ids.remove(slide_ids[0])

        for page_number in range(pdf.page_count):
            page = pdf.load_page(page_number)

            # Render at 2x for good visual quality.
            pixmap = page.get_pixmap(
                matrix=fitz.Matrix(2.0, 2.0),
                alpha=False,
            )

            image_path = os.path.join(
                temp_dir,
                f"page_{page_number + 1}.png",
            )
            pixmap.save(image_path)

            slide = prs.slides.add_slide(
                prs.slide_layouts[6]
            )

            # Fit the page inside the slide while preserving
            # the original PDF page aspect ratio.
            page_width = float(page.rect.width)
            page_height = float(page.rect.height)
            page_ratio = page_width / page_height
            slide_ratio = 13.333333 / 7.5

            if page_ratio >= slide_ratio:
                width = Inches(13.333333)
                height = Inches(13.333333 / page_ratio)
                left = 0
                top = Inches((7.5 - float(height.inches)) / 2)
            else:
                height = Inches(7.5)
                width = Inches(7.5 * page_ratio)
                top = 0
                left = Inches((13.333333 - float(width.inches)) / 2)

            slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=width,
                height=height,
            )

        pdf.close()
        prs.save(pptx_path)

        with open(pptx_path, "rb") as pptx_file:
            output = pptx_file.read()

        if not output:
            raise RuntimeError(
                "The generated PowerPoint file is empty."
            )

        return output, pptx_filename

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True,
        )


# ============================================================
# MAIN CONVERSION API
# ============================================================

@app.post("/api/convert")
async def convert(
    toolId: Annotated[str, Form(...)],
    files: Annotated[UploadFile, File(...)],
):

    if not files:
        return PlainTextResponse(
            "No file was uploaded.",
            status_code=400,
        )

    if not files.filename:
        return PlainTextResponse(
            "The uploaded file has no filename.",
            status_code=400,
        )

    filename = files.filename

    # ========================================================
    # PDF → POWERPOINT
    # ========================================================

    if toolId == "pdf-to-powerpoint":
        if not filename.lower().endswith(".pdf"):
            return PlainTextResponse(
                "PDF to PowerPoint requires a PDF file.",
                status_code=400,
            )

        pdf_bytes = await files.read()

        if not pdf_bytes:
            return PlainTextResponse(
                "The uploaded PDF is empty.",
                status_code=400,
            )

        try:
            print(
                f"PDF to PowerPoint: converting {filename} "
                f"({len(pdf_bytes)} bytes)"
            )

            output, output_filename = convert_pdf_to_powerpoint(
                pdf_bytes,
                filename,
            )

            print(
                f"PDF to PowerPoint: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                headers={
                    "X-Output-Filename": output_filename,
                    "Content-Disposition": (
                        f'attachment; filename="{output_filename}"'
                    ),
                },
            )

        except Exception as error:
            print(
                "PDF to PowerPoint error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this PDF to PowerPoint. "
                f"Error: {error}",
                status_code=500,
            )

    # ========================================================
    # EXCEL → PDF
    # ========================================================

    if toolId == "excel-to-pdf":
        if not filename.lower().endswith((".xlsx", ".xls")):
            return PlainTextResponse(
                "Excel to PDF requires an XLSX or XLS file.",
                status_code=400,
            )

        excel_bytes = await files.read()

        if not excel_bytes:
            return PlainTextResponse(
                "The uploaded Excel file is empty.",
                status_code=400,
            )

        try:
            print(
                f"Excel to PDF: converting {filename} "
                f"({len(excel_bytes)} bytes)"
            )

            (
                output,
                output_filename,
            ) = convert_excel_to_pdf(
                excel_bytes,
                filename,
            )

            print(
                f"Excel to PDF: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type="application/pdf",
                headers={
                    "X-Output-Filename": output_filename,
                    "Content-Disposition": (
                        f'attachment; '
                        f'filename="{output_filename}"'
                    ),
                },
            )

        except Exception as error:
            print(
                "Excel to PDF error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this Excel file to PDF. "
                f"Error: {error}",
                status_code=500,
            )

    # ========================================================
    # WORD → PDF
    # ========================================================

    if toolId == "word-to-pdf":
        if not filename.lower().endswith((".docx", ".doc")):
            return PlainTextResponse(
                "Word to PDF requires a DOCX or DOC file.",
                status_code=400,
            )

        word_bytes = await files.read()

        if not word_bytes:
            return PlainTextResponse(
                "The uploaded Word file is empty.",
                status_code=400,
            )

        try:
            print(
                f"Word to PDF: converting {filename} "
                f"({len(word_bytes)} bytes)"
            )

            (
                output,
                output_filename,
            ) = convert_word_to_pdf(
                word_bytes,
                filename,
            )

            print(
                f"Word to PDF: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type="application/pdf",
                headers={
                    "X-Output-Filename":
                        output_filename,
                    "Content-Disposition":
                        (
                            f'attachment; '
                            f'filename="{output_filename}"'
                        ),
                },
            )

        except Exception as error:
            print(
                "Word to PDF error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this Word file to PDF. "
                f"Error: {error}",
                status_code=500,
            )

    # Existing PDF tools continue below.
    if not filename.lower().endswith(".pdf"):
        return PlainTextResponse(
            "This tool requires a PDF file.",
            status_code=400,
        )

    pdf_bytes = await files.read()

    if not pdf_bytes:
        return PlainTextResponse(
            "The uploaded PDF is empty.",
            status_code=400,
        )

    # ========================================================
    # PDF → EXCEL
    # ========================================================

    if toolId == "pdf-to-excel":

        try:
            print(
                f"PDF to Excel: converting {filename} "
                f"({len(pdf_bytes)} bytes)"
            )

            (
                output,
                output_filename,
            ) = convert_pdf_to_excel(
                pdf_bytes,
                filename,
            )

            print(
                f"PDF to Excel: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "spreadsheetml.sheet"
                ),
                headers={
                    "X-Output-Filename":
                        output_filename,
                    "Content-Disposition":
                        (
                            f'attachment; '
                            f'filename="{output_filename}"'
                        ),
                },
            )

        except Exception as error:

            print(
                "PDF to Excel error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this PDF to Excel. "
                f"Error: {error}",
                status_code=500,
            )

    # ========================================================
    # PDF → WORD
    # ========================================================

    if toolId == "pdf-to-word":

        try:
            print(
                f"PDF to Word: converting {filename} "
                f"({len(pdf_bytes)} bytes)"
            )

            output = convert_pdf_to_docx(
                pdf_bytes
            )

            output_filename = (
                f"{Path(filename).stem}.docx"
            )

            print(
                f"PDF to Word: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                ),
                headers={
                    "X-Output-Filename":
                        output_filename,
                    "Content-Disposition":
                        (
                            f'attachment; '
                            f'filename="{output_filename}"'
                        ),
                },
            )

        except Exception as error:

            print(
                "PDF to Word error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this PDF to Word. "
                f"Error: {error}",
                status_code=500,
            )

    # ========================================================
    # PDF → JPG
    # ========================================================

    if toolId == "pdf-to-jpg":

        try:
            print(
                f"PDF to JPG: converting {filename} "
                f"({len(pdf_bytes)} bytes)"
            )

            (
                output,
                media_type,
                output_filename,
            ) = convert_pdf_to_jpg(
                pdf_bytes,
                filename,
            )

            print(
                f"PDF to JPG: success - "
                f"{output_filename} "
                f"({len(output)} bytes)"
            )

            return StreamingResponse(
                iter([output]),
                media_type=media_type,
                headers={
                    "X-Output-Filename":
                        output_filename,
                    "Content-Disposition":
                        (
                            f'attachment; '
                            f'filename="{output_filename}"'
                        ),
                },
            )

        except Exception as error:

            print(
                "PDF to JPG error:",
                repr(error),
            )

            return PlainTextResponse(
                "Unable to convert this PDF to JPG. "
                f"Error: {error}",
                status_code=500,
            )

    # ========================================================
    # OTHER TOOLS
    # ========================================================

    return PlainTextResponse(
        f"Conversion tool '{toolId}' is not implemented "
        "in the Python conversion server.",
        status_code=501,
    )


# ============================================================
# START SERVER
# ============================================================

if __name__ == "__main__":
    import uvicorn

    uvicorn.run(
        app,
        host="127.0.0.1",
        port=8000,
    )